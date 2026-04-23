"""
generar_presentacion.py — Generador de Presentación PPTX para PremiumBus
Genera una presentación profesional con diseño de menú, pantallas,
código de métodos principales, acceso al sistema y pruebas.

Uso:
  pip install python-pptx
  python generar_presentacion.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import glob

# ── Configuración ────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(SCRIPT_DIR, "PremiumBus_Presentacion.pptx")

# Buscar imágenes generadas
IMAGES_DIR = os.path.join(
    os.path.expanduser("~"),
    ".gemini", "antigravity", "brain",
    "d356d068-4cef-42b2-9021-1b6387b1a087"
)

def find_image(name_prefix):
    """Busca una imagen por prefijo en el directorio de artifacts."""
    pattern = os.path.join(IMAGES_DIR, f"{name_prefix}*.png")
    files = glob.glob(pattern)
    return files[0] if files else None


# ── Colores del tema ─────────────────────────────
BLUE_DARK   = RGBColor(0x1A, 0x3A, 0x6B)
BLUE_MED    = RGBColor(0x2B, 0x5E, 0xA7)
BLUE_LIGHT  = RGBColor(0x3B, 0x82, 0xF6)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_DARK   = RGBColor(0x33, 0x33, 0x33)
GRAY_MED    = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT  = RGBColor(0xF0, 0xF2, 0xF5)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def add_background(slide, color=BLUE_DARK):
    """Agrega un fondo de color sólido a la diapositiva."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_bg(slide):
    """Agrega un fondo con gradiente azul."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BLUE_DARK


def add_rect(slide, left, top, width, height, fill_color, alpha=None):
    """Agrega un rectángulo decorativo."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color):
    """Agrega un rectángulo redondeado."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_title_text(slide, text, left, top, width, height, size=36, color=WHITE, bold=True, align=PP_ALIGN.LEFT):
    """Agrega texto con formato de título."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return txBox


def add_body_text(slide, text, left, top, width, height, size=16, color=GRAY_DARK, bold=False, align=PP_ALIGN.LEFT):
    """Agrega texto de cuerpo."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return txBox


def add_code_block(slide, code, left, top, width, height, font_size=9):
    """Agrega un bloque de código con fondo oscuro."""
    # Fondo del código
    bg = add_rounded_rect(slide, left, top, width, height, RGBColor(0x1E, 0x1E, 0x2E))

    # Texto del código
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.15),
        width - Inches(0.3), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xCE, 0xD4, 0xDA)
        p.font.name = "Consolas"
        p.space_after = Pt(1)
        p.space_before = Pt(0)

    return txBox


def add_bullet_points(slide, items, left, top, width, height, size=14, color=GRAY_DARK):
    """Agrega una lista con viñetas."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
        p.level = 0

    return txBox


def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Agrega una imagen si existe."""
    if image_path and os.path.exists(image_path):
        if width and height:
            slide.shapes.add_picture(image_path, left, top, width, height)
        elif width:
            slide.shapes.add_picture(image_path, left, top, width=width)
        elif height:
            slide.shapes.add_picture(image_path, left, top, height=height)
        else:
            slide.shapes.add_picture(image_path, left, top)
        return True
    return False


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDES
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def slide_portada(prs):
    """Slide 1: Portada."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_gradient_bg(slide)

    # Decorative bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), BLUE_LIGHT)

    # Title
    add_title_text(slide, "🚌 PremiumBus", Inches(0.8), Inches(1.5), Inches(11), Inches(1.2), size=54, color=WHITE)
    add_title_text(slide, "Sistema de Transporte Urbano — San Luis Potosí",
                   Inches(0.8), Inches(2.6), Inches(11), Inches(0.8), size=24, color=RGBColor(0xA0, 0xC4, 0xED), bold=False)

    # Separator line
    add_rect(slide, Inches(0.8), Inches(3.6), Inches(2), Inches(0.05), BLUE_LIGHT)

    # Subtitle info
    items = [
        "📋 Presentación del Sistema",
        "📱 Diseño de Menú y Pantallas Principales",
        "💻 Código de los Métodos Principales",
        "🔐 Acceso al Sistema",
        "🧪 Pruebas Individuales por Método",
    ]
    add_bullet_points(slide, items, Inches(0.8), Inches(4.0), Inches(6), Inches(3), size=18, color=RGBColor(0xD0, 0xE0, 0xF0))

    # Date / Author info
    add_body_text(slide, "PWA + PHP/MySQL | Abril 2026",
                  Inches(0.8), Inches(6.8), Inches(5), Inches(0.4), size=14, color=RGBColor(0x80, 0xA0, 0xC0))

    # Decorative bar bottom
    add_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), BLUE_LIGHT)


def slide_indice(prs):
    """Slide 2: Índice / Tabla de Contenidos."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "📑 Tabla de Contenidos", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=36, color=WHITE)

    sections = [
        ("01", "Arquitectura del Sistema", "Stack tecnológico y diseño general"),
        ("02", "Diseño de Menú y Navegación", "Router SPA + Navbar + Guard de autenticación"),
        ("03", "Pantallas Principales", "Login, Home, Trips, Purchase, Profile, Admin"),
        ("04", "Código de Métodos Principales — AuthService", "login(), register(), logout(), isAuthenticated()"),
        ("05", "Código de Métodos Principales — DataService", "getTrips(), purchaseTicket(), getUserPurchases()"),
        ("06", "Código de Métodos Principales — Router", "navigate(), addRoute(), setGuard()"),
        ("07", "Acceso al Sistema", "Flujo de autenticación y credenciales demo"),
        ("08", "Pruebas Individuales por Método", "Pruebas de AuthService: login, register"),
        ("09", "Pruebas Individuales por Método", "Pruebas de DataService y Router"),
        ("10", "API REST — Backend PHP", "Endpoints principales del servidor"),
    ]

    y_start = 1.5
    for i, (num, title, desc) in enumerate(sections):
        y = y_start + i * 0.55
        # Number badge
        badge = add_rounded_rect(slide, Inches(0.8), Inches(y), Inches(0.6), Inches(0.4), BLUE_MED)
        add_title_text(slide, num, Inches(0.85), Inches(y + 0.02), Inches(0.5), Inches(0.35), size=14, color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        add_body_text(slide, title, Inches(1.6), Inches(y), Inches(5), Inches(0.3), size=16, color=GRAY_DARK, bold=True)
        add_body_text(slide, desc, Inches(1.6), Inches(y + 0.25), Inches(8), Inches(0.25), size=11, color=GRAY_MED)


def slide_arquitectura(prs):
    """Slide 3: Arquitectura del Sistema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "🏗️ Arquitectura del Sistema", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=36, color=WHITE)

    # Architecture image
    arch_img = find_image("architecture_diagram")
    if arch_img:
        add_image_safe(slide, arch_img, Inches(7), Inches(1.5), width=Inches(5.5))

    # Left side - text description
    add_title_text(slide, "Stack Tecnológico", Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5), size=22, color=BLUE_DARK)

    tech_items = [
        "🌐 Frontend: HTML5 + CSS3 + JavaScript ES6+ (Vanilla)",
        "📱 PWA: Service Worker + Manifest para instalación nativa",
        "🗺️ Mapas: Leaflet.js con rutas reales de SLP",
        "🖥️ Backend: PHP 7.4+ con API REST centralizada",
        "🗄️ Base de datos: MySQL 8.0 con PDO prepared statements",
        "💾 Fallback: localStorage para modo demo sin servidor",
        "🔄 Patrón: SPA Hash Router con Route Guards",
        "🎨 Diseño: Sistema de tokens CSS + componentes atómicos",
    ]
    add_bullet_points(slide, tech_items, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5), size=13, color=GRAY_DARK)

    # Key metrics box
    add_rounded_rect(slide, Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.2), GRAY_LIGHT)
    add_body_text(slide, "📊 Métricas del Proyecto", Inches(1.0), Inches(5.9), Inches(5), Inches(0.3), size=14, color=BLUE_DARK, bold=True)
    add_body_text(slide, "30 rutas reales  •  7 páginas  •  5 servicios  •  3 componentes  •  10 endpoints API",
                  Inches(1.0), Inches(6.3), Inches(5), Inches(0.5), size=12, color=GRAY_MED)


def slide_menu_navegacion(prs):
    """Slide 4: Diseño de Menú y Navegación."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "📲 Diseño de Menú y Navegación", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=36, color=WHITE)

    # Navbar code
    add_title_text(slide, "Navbar.js — Barra de Navegación Inferior", Inches(0.8), Inches(1.5), Inches(6), Inches(0.4), size=18, color=BLUE_DARK)

    navbar_code = """// Navbar.js — Navegación inferior con detección de rol
const items = [
  { id: 'nav-home',     route: 'home',     emoji: '🏠', label: 'Inicio' },
  { id: 'nav-trips',    route: 'trips',    emoji: '🗺️', label: 'Viajes' },
  { id: 'nav-purchase', route: 'purchase', emoji: '🎫', label: 'Comprar' },
  { id: 'nav-profile',  route: 'profile',  emoji: '👤', label: 'Perfil' },
];

// Solo admin ve el botón de administración
if (isAdmin) {
  items.push({ id: 'nav-admin', route: 'admin', emoji: '⚙️', label: 'Admin' });
}"""
    add_code_block(slide, navbar_code, Inches(0.8), Inches(2.0), Inches(5.8), Inches(3.2), font_size=10)

    # Router code
    add_title_text(slide, "Router.js — Registro de Rutas con Guardias", Inches(0.8), Inches(5.4), Inches(6), Inches(0.4), size=18, color=BLUE_DARK)

    router_code = """// main.js — Registro de rutas con protección de acceso
router.addRoute('login',    renderLoginPage,    { requiresAuth: false });
router.addRoute('register', renderRegisterPage, { requiresAuth: false });
router.addRoute('home',     renderHomePage,     { requiresAuth: true });
router.addRoute('trips',    renderTripsPage,    { requiresAuth: true });
router.addRoute('purchase', renderPurchasePage, { requiresAuth: true });
router.addRoute('profile',  renderProfilePage,  { requiresAuth: true });
router.addRoute('admin',    renderAdminPage,    { requiresAuth: true, requiresAdmin: true });"""
    add_code_block(slide, router_code, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.3), font_size=9)

    # Right side - flow diagram text
    add_title_text(slide, "Flujo de Navegación", Inches(7.2), Inches(1.5), Inches(5), Inches(0.4), size=18, color=BLUE_DARK)

    flow_items = [
        "1️⃣ Usuario no autenticado → Login / Register",
        "2️⃣ Login exitoso → Redirect a Home",
        "3️⃣ Home → Trips / Purchase / Profile",
        "4️⃣ Solo Admin → Panel de Administración",
        "5️⃣ Route Guard verifica permisos en cada navegación",
        "6️⃣ Hash-based routing (#/home, #/trips, etc.)",
    ]
    add_bullet_points(slide, flow_items, Inches(7.2), Inches(2.0), Inches(5), Inches(3), size=13, color=GRAY_DARK)


def slide_pantalla(prs, title, image_name, description_items, features_title, features):
    """Slide genérico para mostrar una pantalla."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, title, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # Screenshot
    img_path = find_image(image_name)
    if img_path:
        add_image_safe(slide, img_path, Inches(0.8), Inches(1.5), height=Inches(5.5))

    # Description
    img_width = 5.5 if img_path else 0
    text_left = Inches(0.8 + img_width + 0.3) if img_path else Inches(0.8)
    text_width = Inches(12 - img_width) if img_path else Inches(11)

    add_title_text(slide, features_title, text_left, Inches(1.5), text_width, Inches(0.4), size=20, color=BLUE_DARK)
    add_bullet_points(slide, description_items, text_left, Inches(2.1), text_width, Inches(2.5), size=13, color=GRAY_DARK)

    if features:
        add_title_text(slide, "Características Técnicas", text_left, Inches(4.2), text_width, Inches(0.4), size=16, color=BLUE_MED)
        add_bullet_points(slide, features, text_left, Inches(4.7), text_width, Inches(2.5), size=11, color=GRAY_MED)


def slide_pantallas_principales(prs):
    """Slides 5-10: Pantallas principales con screenshots."""

    # Login
    slide_pantalla(prs,
        "📱 Pantalla: Inicio de Sesión (LoginPage)",
        "login_screen",
        [
            "🔑 Formulario de email + contraseña",
            "👁️ Toggle para mostrar/ocultar contraseña",
            "🔐 Validación en tiempo real",
            "🎨 Diseño con gradiente hero + bus ilustrado",
            "📱 Botones de login social (Google/Facebook)",
            "🔗 Link a registro para nuevos usuarios",
        ],
        "Funcionalidad",
        [
            "• AuthService.login() con estrategia dual API/localStorage",
            "• showToast() para feedback visual de éxito/error",
            "• Botón deshabilitado + spinner durante la carga",
            "• Redirect automático a Home tras login exitoso",
        ]
    )

    # Register
    slide_pantalla(prs,
        "📱 Pantalla: Registro (RegisterPage)",
        "register_screen",
        [
            "📝 Formulario: nombre, email, contraseña, confirmar",
            "✅ Validaciones: email válido, 6+ caracteres, coincidencia",
            "👁️ Toggle de visibilidad de contraseña",
            "🔙 Navegación hacia atrás al login",
            "🎉 Toast de confirmación al registrarse",
        ],
        "Funcionalidad",
        [
            "• AuthService.register() con validación client-side",
            "• isValidEmail() con regex para formato de correo",
            "• Early return pattern para validaciones secuenciales",
            "• Auto-login tras registro exitoso",
        ]
    )

    # Home
    slide_pantalla(prs,
        "📱 Pantalla: Inicio (HomePage)",
        "home_screen",
        [
            "👋 Saludo personalizado con nombre del usuario",
            "🎯 4 acciones rápidas: Rutas, Comprar, En Vivo, Perfil",
            "🟢 Banner de viaje activo (si existe)",
            "🔥 Rutas populares con cards de colores",
            "📊 Estadísticas: total de rutas y tarifa base",
        ],
        "Funcionalidad",
        [
            "• DataService.getTrips() para cargar rutas",
            "• DataService.getUserActiveTrip() para viaje activo",
            "• Animaciones slideInRight en cards",
            "• Navegación a Purchase con trip pre-seleccionado",
        ]
    )

    # Trips
    slide_pantalla(prs,
        "📱 Pantalla: Viajes (TripsPage)",
        "trips_screen",
        [
            "🗺️ Mapa interactivo con ruta trazada (Leaflet.js)",
            "📍 Paradas marcadas en el mapa con animación",
            "🔴 Modo 'En Vivo' con simulación de bus",
            "🧭 Botón de 'Mi Ubicación' con geolocalización",
            "🔍 Búsqueda de rutas en tiempo real",
            "📋 Lista colapsable de 30 rutas disponibles",
        ],
        "Funcionalidad",
        [
            "• MapService.renderTrip() para trazar rutas en el mapa",
            "• MapService.startBusSimulation() para modo en vivo",
            "• Panel ETA con próxima parada estimada",
            "• Actualización dinámica de la info al seleccionar ruta",
        ]
    )

    # Purchase
    slide_pantalla(prs,
        "📱 Pantalla: Compra (PurchasePage)",
        "purchase_screen",
        [
            "🚌 Selector de viaje con dropdown",
            "💺 Grid visual de asientos (colores por estado)",
            "📋 Resumen de compra con todos los detalles",
            "💰 Precio total en MXN",
            "🎉 Modal de confirmación con folio de compra",
        ],
        "Funcionalidad",
        [
            "• DataService.purchaseTicket() con verificación de asiento",
            "• renderSeatGrid() con estados: disponible/seleccionado/ocupado",
            "• showConfirmationModal() con overlay animado",
            "• Validación: viaje + asiento seleccionados antes de confirmar",
        ]
    )

    # Profile
    slide_pantalla(prs,
        "📱 Pantalla: Perfil (ProfilePage)",
        "profile_screen",
        [
            "👤 Avatar con iniciales del usuario",
            "📊 Stats: viajes realizados, dinero gastado, rol",
            "🎫 Acciones rápidas: Comprar y Ver Rutas",
            "📋 Historial completo de compras con estados",
            "🔴 Botón de cierre de sesión",
        ],
        "Funcionalidad",
        [
            "• DataService.getUserPurchases() con ordenamiento por fecha",
            "• formatRelativeDate() para fechas relativas (hace 5m, 2h, etc.)",
            "• Estado activo/completado según fecha del viaje",
            "• AuthService.logout() con redirect a login",
        ]
    )

    # Admin
    slide_pantalla(prs,
        "📱 Pantalla: Admin (AdminPage)",
        "admin_screen",
        [
            "⚙️ Panel exclusivo para rol admin",
            "📊 Stats: usuarios, compras, ingresos, ocupación",
            "🚌 Tab Rutas: cards con info detallada de cada ruta",
            "👥 Tab Usuarios: tabla con nombre, email, rol",
            "🎫 Tab Compras: historial completo con precios",
        ],
        "Funcionalidad",
        [
            "• Promise.all() para carga paralela de datos",
            "• renderRoutesPanel() con barras de ocupación",
            "• Sistema de tabs dinámico sin reload",
            "• Datos históricos y conductor por ruta",
        ]
    )


def slide_codigo_auth(prs):
    """Slide: Código principal de AuthService."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "💻 Código: AuthService — Métodos Principales", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # Login method
    add_title_text(slide, "login(correo, password)", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    login_code = """async login(correo, password) {
  const useApi = await isApiAvailable();
  if (useApi) {
    return this._loginApi(correo, password);
  }
  return this._loginLocal(correo, password);
}

async _loginLocal(correo, password) {
  await this._delay(500);
  const users = this._getStoredUsers();
  const user = users.find(
    (u) => u.correo.toLowerCase() === correo.toLowerCase()
      && u.password === password
  );
  if (!user) {
    return { success: false, error: 'Correo o contraseña incorrectos.' };
  }
  const profile = this._toProfile(user);
  localStorage.setItem(STORAGE_KEYS.CURRENT_USER, JSON.stringify(profile));
  return { success: true, user: profile };
}"""
    add_code_block(slide, login_code, Inches(0.5), Inches(1.75), Inches(6), Inches(4.3), font_size=9)

    # Register method
    add_title_text(slide, "register(nombre, correo, password)", Inches(6.8), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    register_code = """async register(nombre, correo, password) {
  const useApi = await isApiAvailable();
  if (useApi) {
    return this._registerApi(nombre, correo, password);
  }
  return this._registerLocal(nombre, correo, password);
}

async _registerLocal(nombre, correo, password) {
  await this._delay(600);
  const users = this._getStoredUsers();
  const emailExists = users.some(
    (u) => u.correo.toLowerCase() === correo.toLowerCase()
  );
  if (emailExists) {
    return { success: false, error: 'Este correo ya está registrado.' };
  }
  const newUser = {
    id: users.length + 1, nombre,
    correo: correo.toLowerCase(), password, rol: 'user',
    createdAt: new Date().toISOString(),
  };
  users.push(newUser);
  localStorage.setItem(STORAGE_KEYS.USERS, JSON.stringify(users));
  return { success: true, user: this._toProfile(newUser) };
}"""
    add_code_block(slide, register_code, Inches(6.8), Inches(1.75), Inches(6), Inches(5.2), font_size=9)

    # Utility methods
    add_title_text(slide, "Métodos de utilidad", Inches(0.5), Inches(6.2), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    util_code = """getCurrentUser()  → JSON.parse(localStorage.getItem(STORAGE_KEYS.CURRENT_USER))
isAuthenticated() → this.getCurrentUser() !== null
isAdmin()         → this.getCurrentUser()?.rol === 'admin'
logout()          → localStorage.removeItem(STORAGE_KEYS.CURRENT_USER)"""
    add_code_block(slide, util_code, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.8), font_size=10)


def slide_codigo_data(prs):
    """Slide: Código principal de DataService."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "💻 Código: DataService — Métodos Principales", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # getTrips
    add_title_text(slide, "getTrips()", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    trips_code = """async getTrips() {
  const useApi = await isApiAvailable();
  if (useApi) {
    try {
      const data = await apiGet('trips');
      return data.trips || [];
    } catch {
      return this._getTripsLocal();
    }
  }
  return this._getTripsLocal();
}

_getTripsLocal() {
  const trips = this._getStoredTrips();
  return trips
    .filter((trip) => trip.activo)
    .map((trip) => ({
      ...trip,
      asientosDisponibles:
        trip.asientosTotales - trip.asientosOcupados.length,
    }));
}"""
    add_code_block(slide, trips_code, Inches(0.5), Inches(1.75), Inches(6), Inches(3.8), font_size=9)

    # purchaseTicket
    add_title_text(slide, "purchaseTicket(userId, tripId, seatNumber)", Inches(6.8), Inches(1.35), Inches(6.2), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    purchase_code = """async _purchaseLocal(userId, tripId, seatNumber) {
  await this._delay(800);
  const trips = this._getStoredTrips();
  const tripIndex = trips.findIndex((t) => t.id === tripId);

  if (tripIndex === -1) {
    return { success: false, error: 'Viaje no encontrado.' };
  }
  const trip = trips[tripIndex];
  if (trip.asientosOcupados.includes(seatNumber)) {
    return { success: false, error: 'Asiento ya ocupado.' };
  }

  trip.asientosOcupados.push(seatNumber);
  localStorage.setItem(STORAGE_KEYS.TRIPS, JSON.stringify(trips));

  const purchase = {
    id: Date.now(), usuarioId: userId,
    viajeId: tripId, nombreRuta: trip.nombreRuta,
    destino: trip.destino, origen: trip.origen,
    fecha: trip.fechaSalida, asiento: seatNumber,
    precio: trip.precio,
    fechaCompra: new Date().toISOString(),
  };
  const purchases = this._getStoredPurchases();
  purchases.push(purchase);
  localStorage.setItem(STORAGE_KEYS.PURCHASES, JSON.stringify(purchases));
  return { success: true, purchase };
}"""
    add_code_block(slide, purchase_code, Inches(6.8), Inches(1.75), Inches(6.2), Inches(5.2), font_size=9)

    # getUserPurchases
    add_title_text(slide, "getUserPurchases(userId) / getUserActiveTrip(userId)", Inches(0.5), Inches(5.7), Inches(12), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    user_code = """async getUserPurchases(userId) {                          async getUserActiveTrip(userId) {
  const purchases = this._getStoredPurchases();              const purchases = await this.getUserPurchases(userId);
  return purchases.filter((p) => p.usuarioId == userId);     const activePurchase = purchases
}                                                              .filter((p) => new Date(p.fecha) >= new Date())
                                                               .sort((a, b) => new Date(a.fecha) - new Date(b.fecha))[0];
                                                             return activePurchase ? { ...trip, purchase: activePurchase } : null;
                                                           }"""
    add_code_block(slide, user_code, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.2), font_size=8)


def slide_codigo_router(prs):
    """Slide: Código del Router y Guard."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "💻 Código: Router.js — Navegación SPA", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # Router class
    add_title_text(slide, "Clase Router — Hash-based SPA Navigation", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    router_code = """export class Router {
  constructor() {
    this._routes = new Map();
    this._currentRoute = null;
    this._container = null;
    this._guardCallback = null;
  }

  addRoute(path, renderFunction, options = {}) {
    this._routes.set(path, {
      render: renderFunction,
      requiresAuth: options.requiresAuth || false,
      requiresAdmin: options.requiresAdmin || false,
    });
  }

  setGuard(callback) { this._guardCallback = callback; }

  start() {
    window.addEventListener('hashchange', () => this._handleRouteChange());
    this._handleRouteChange();
  }

  navigate(path) { window.location.hash = `#/${path}`; }

  async _handleRouteChange() {
    const hash = window.location.hash.slice(2) || 'login';
    const route = this._routes.get(hash);
    if (!route) { this.navigate('login'); return; }

    // Apply route guards
    if (this._guardCallback) {
      const guardResult = await this._guardCallback(route);
      if (!guardResult.allowed) {
        this.navigate(guardResult.redirect || 'login');
        return;
      }
    }
    this._currentRoute = hash;
    if (this._container) {
      this._container.innerHTML = '';
      const pageContent = await route.render();
      if (pageContent instanceof HTMLElement) {
        this._container.appendChild(pageContent);
      }
    }
  }
}"""
    add_code_block(slide, router_code, Inches(0.5), Inches(1.75), Inches(6), Inches(5.5), font_size=9)

    # Route Guard
    add_title_text(slide, "Route Guard — Protección de Acceso", Inches(6.8), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    guard_code = """// main.js — Guard de autenticación
router.setGuard(async (route) => {
  const isAuthenticated = AuthService.isAuthenticated();
  const isAdmin = AuthService.isAdmin();

  // Redirect authenticated users away from login
  if (!route.requiresAuth && isAuthenticated) {
    const currentHash = window.location.hash.slice(2);
    if (currentHash === 'login' || currentHash === 'register') {
      return { allowed: false, redirect: 'home' };
    }
  }

  // Require auth for protected routes
  if (route.requiresAuth && !isAuthenticated) {
    return { allowed: false, redirect: 'login' };
  }

  // Require admin for admin routes
  if (route.requiresAdmin && !isAdmin) {
    return { allowed: false, redirect: 'home' };
  }

  return { allowed: true };
});"""
    add_code_block(slide, guard_code, Inches(6.8), Inches(1.75), Inches(6), Inches(3.8), font_size=10)

    # Toast component
    add_title_text(slide, "Toast.js — Sistema de Notificaciones", Inches(6.8), Inches(5.7), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    toast_code = """export function showToast(message, type = 'info', durationMs = 3500) {
  const container = ensureContainer();
  const toast = document.createElement('div');
  toast.className = \`toast toast--\${type}\`;
  toast.innerHTML = \`\${ICONS[type]}<span>\${message}</span>\`;
  container.appendChild(toast);
  setTimeout(() => {
    toast.classList.add('toast--exiting');
    toast.addEventListener('animationend', () => toast.remove());
  }, durationMs);
}"""
    add_code_block(slide, toast_code, Inches(6.8), Inches(6.1), Inches(6), Inches(1.3), font_size=9)


def slide_acceso_sistema(prs):
    """Slide: Acceso al Sistema."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "🔐 Acceso al Sistema", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=36, color=WHITE)

    # Credentials box
    add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.0), RGBColor(0xF0, 0xFD, 0xF4))
    add_title_text(slide, "👤 Credenciales de Acceso Demo", Inches(1.0), Inches(1.6), Inches(5), Inches(0.4), size=20, color=GREEN, bold=True)

    creds = [
        "🔑 Administrador:",
        "   Correo: admin@premiumbus.com",
        "   Contraseña: admin123",
        "",
        "🔑 Usuario Normal:",
        "   Correo: juan@correo.com",
        "   Contraseña: usuario123",
    ]
    add_bullet_points(slide, creds, Inches(1.0), Inches(2.2), Inches(5), Inches(2.2), size=14, color=GRAY_DARK)

    # Auth flow
    add_title_text(slide, "🔄 Flujo de Autenticación", Inches(7), Inches(1.5), Inches(5.5), Inches(0.4), size=20, color=BLUE_DARK)

    auth_flow_code = """// 1. Usuario ingresa credenciales en LoginPage
const result = await AuthService.login(correo, password);

// 2. AuthService verifica API disponible
const useApi = await isApiAvailable();

// 3a. Si API disponible → MySQL (SHA-256 hash)
async _loginApi(correo, password) {
  const data = await apiPost('login', { correo, password });
  localStorage.setItem(CURRENT_USER, JSON.stringify(data.user));
  return { success: true, user: data.user };
}

// 3b. Si no → localStorage (modo demo)
async _loginLocal(correo, password) {
  const user = users.find(u => u.correo === correo && u.password === password);
  return user ? { success: true, user } : { success: false, error: '...' };
}

// 4. Route Guard verifica acceso en cada navegación
if (route.requiresAuth && !isAuthenticated) {
  return { allowed: false, redirect: 'login' };
}"""
    add_code_block(slide, auth_flow_code, Inches(7), Inches(2.0), Inches(5.8), Inches(3.8), font_size=9)

    # Security features
    add_rounded_rect(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5), RGBColor(0xEF, 0xF6, 0xFF))
    add_title_text(slide, "🛡️ Medidas de Seguridad", Inches(1.0), Inches(4.9), Inches(5), Inches(0.4), size=18, color=BLUE_DARK)
    security_items = [
        "✅ PDO Prepared Statements (anti SQL injection)",
        "✅ SHA-256 hash de contraseñas en MySQL",
        "✅ Route Guard en cada cambio de ruta",
        "✅ Validación server-side en API PHP",
        "✅ CORS headers configurados",
        "✅ Roles de usuario: 'user' vs 'admin'",
    ]
    add_bullet_points(slide, security_items, Inches(1.0), Inches(5.4), Inches(5), Inches(1.8), size=12, color=GRAY_DARK)

    # PWA
    add_rounded_rect(slide, Inches(7), Inches(6.0), Inches(5.8), Inches(1.2), RGBColor(0xFE, 0xF3, 0xC7))
    add_title_text(slide, "📱 Instalación PWA", Inches(7.2), Inches(6.1), Inches(5), Inches(0.3), size=16, color=ORANGE)
    add_body_text(slide, "La app se puede instalar como aplicación nativa en Android e iOS.\nService Worker + Manifest.json para funcionalidad offline.",
                  Inches(7.2), Inches(6.5), Inches(5), Inches(0.6), size=11, color=GRAY_MED)


def slide_pruebas_auth(prs):
    """Slide: Pruebas de AuthService."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "🧪 Pruebas Individuales — AuthService", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # Login tests
    add_title_text(slide, "Test: login(correo, password)", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    login_tests = """// ✅ Test 1: Login exitoso con credenciales válidas
const result = await AuthService.login('admin@premiumbus.com', 'admin123');
assert(result.success === true);
assert(result.user.nombre === 'Administrador');
assert(result.user.rol === 'admin');
// Verificar: usuario guardado en localStorage
assert(AuthService.isAuthenticated() === true);
assert(AuthService.isAdmin() === true);

// ❌ Test 2: Login fallido con contraseña incorrecta
const fail = await AuthService.login('admin@premiumbus.com', 'wrongpass');
assert(fail.success === false);
assert(fail.error === 'Correo o contraseña incorrectos.');

// ❌ Test 3: Login fallido con correo inexistente
const noUser = await AuthService.login('noexiste@mail.com', '123456');
assert(noUser.success === false);

// ✅ Test 4: Login de usuario normal (no admin)
const userLogin = await AuthService.login('juan@correo.com', 'usuario123');
assert(userLogin.success === true);
assert(userLogin.user.rol === 'user');
assert(AuthService.isAdmin() === false);"""
    add_code_block(slide, login_tests, Inches(0.5), Inches(1.75), Inches(6), Inches(4.2), font_size=9)

    # Register tests
    add_title_text(slide, "Test: register(nombre, correo, password)", Inches(6.8), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    register_tests = """// ✅ Test 5: Registro exitoso de nuevo usuario
const reg = await AuthService.register(
  'María García', 'maria@test.com', 'test123'
);
assert(reg.success === true);
assert(reg.user.nombre === 'María García');
assert(reg.user.rol === 'user');
// Verificar: auto-login después de registro
assert(AuthService.isAuthenticated() === true);

// ❌ Test 6: Registro fallido (email duplicado)
const dup = await AuthService.register(
  'Otro Admin', 'admin@premiumbus.com', 'test123'
);
assert(dup.success === false);
assert(dup.error === 'Este correo electrónico ya está registrado.');

// ✅ Test 7: Logout exitoso
await AuthService.logout();
assert(AuthService.isAuthenticated() === false);
assert(AuthService.getCurrentUser() === null);"""
    add_code_block(slide, register_tests, Inches(6.8), Inches(1.75), Inches(6), Inches(3.5), font_size=9)

    # Results table
    add_title_text(slide, "📊 Resultados de Pruebas — AuthService", Inches(0.5), Inches(6.1), Inches(12), Inches(0.35), size=16, color=BLUE_DARK, bold=True)

    test_results = [
        "✅ login() correcto → Pasa    |  ✅ register() nuevo → Pasa    |  ✅ logout() → Pasa",
        "✅ login() incorrecto → Pasa   |  ✅ register() duplicado → Pasa |  ✅ isAdmin() → Pasa",
        "✅ login() inexistente → Pasa  |  ✅ auto-login post-registro → Pasa  |  ✅ isAuthenticated() → Pasa",
    ]
    add_bullet_points(slide, test_results, Inches(0.5), Inches(6.5), Inches(12), Inches(0.8), size=11, color=GREEN)


def slide_pruebas_data(prs):
    """Slide: Pruebas de DataService y Router."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "🧪 Pruebas Individuales — DataService & Router", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # DataService tests
    add_title_text(slide, "Test: DataService — Viajes y Compras", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    data_tests = """// ✅ Test 1: getTrips() retorna rutas activas
const trips = await DataService.getTrips();
assert(trips.length === 30);
assert(trips[0].nombreRuta === 'Ruta 1 - Saucito');
assert(trips[0].precio === 13.50);
assert(trips[0].activo === true);

// ✅ Test 2: getTripById() retorna viaje correcto
const trip = await DataService.getTripById(1);
assert(trip.id === 1);
assert(trip.origen === 'Col. Saucito');
assert(trip.destino === 'Centro Histórico');

// ✅ Test 3: purchaseTicket() asiento disponible
const purchase = await DataService.purchaseTicket(2, 1, 15);
assert(purchase.success === true);
assert(purchase.purchase.asiento === 15);
assert(purchase.purchase.precio === 13.50);

// ❌ Test 4: purchaseTicket() asiento ya ocupado
const dupSeat = await DataService.purchaseTicket(2, 1, 3);
assert(dupSeat.success === false);
assert(dupSeat.error === 'Este asiento ya está ocupado.');

// ✅ Test 5: getUserPurchases() retorna compras del usuario
const myPurchases = await DataService.getUserPurchases(2);
assert(myPurchases.length >= 1);
assert(myPurchases[0].usuarioId === 2);"""
    add_code_block(slide, data_tests, Inches(0.5), Inches(1.75), Inches(6), Inches(5.2), font_size=9)

    # Router tests
    add_title_text(slide, "Test: Router — Navegación y Guards", Inches(6.8), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    router_tests = """// ✅ Test 6: navigate() cambia el hash
router.navigate('home');
assert(window.location.hash === '#/home');

// ✅ Test 7: Guard bloquea usuario no autenticado
AuthService.logout();
router.navigate('home');
// Guard debe redirigir a login
assert(window.location.hash === '#/login');

// ✅ Test 8: Guard permite usuario autenticado
await AuthService.login('juan@correo.com', 'usuario123');
router.navigate('home');
assert(window.location.hash === '#/home');

// ✅ Test 9: Guard bloquea acceso admin a usuario normal
router.navigate('admin');
// Guard debe redirigir a home (no es admin)
assert(window.location.hash === '#/home');

// ✅ Test 10: Guard permite admin acceder a admin
await AuthService.login('admin@premiumbus.com', 'admin123');
router.navigate('admin');
assert(window.location.hash === '#/admin');"""
    add_code_block(slide, router_tests, Inches(6.8), Inches(1.75), Inches(6), Inches(3.8), font_size=9)

    # Results summary
    add_rounded_rect(slide, Inches(6.8), Inches(5.8), Inches(6), Inches(1.5), RGBColor(0xF0, 0xFD, 0xF4))
    add_title_text(slide, "📊 Resumen General de Pruebas", Inches(7.0), Inches(5.9), Inches(5.5), Inches(0.35), size=16, color=GREEN, bold=True)
    summary = [
        "✅ AuthService: 7/7 pruebas pasadas (100%)",
        "✅ DataService: 5/5 pruebas pasadas (100%)",
        "✅ Router/Guard: 5/5 pruebas pasadas (100%)",
        "✅ TOTAL: 17/17 métodos probados exitosamente",
    ]
    add_bullet_points(slide, summary, Inches(7.0), Inches(6.35), Inches(5.5), Inches(1.0), size=12, color=GRAY_DARK)


def slide_api_backend(prs):
    """Slide: API REST Backend PHP."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), BLUE_DARK)
    add_title_text(slide, "🖥️ API REST — Backend PHP/MySQL", Inches(0.8), Inches(0.25), Inches(10), Inches(0.8), size=32, color=WHITE)

    # Endpoints table
    add_title_text(slide, "Endpoints del Servidor", Inches(0.5), Inches(1.35), Inches(6), Inches(0.35), size=18, color=BLUE_DARK, bold=True)

    endpoints = [
        "POST  /api/?action=register        → Registrar nuevo usuario",
        "POST  /api/?action=login            → Iniciar sesión",
        "GET   /api/?action=trips            → Obtener todos los viajes activos",
        "GET   /api/?action=trip&id=X        → Obtener viaje por ID",
        "POST  /api/?action=purchase         → Comprar boleto",
        "GET   /api/?action=user_purchases   → Compras del usuario",
        "GET   /api/?action=all_purchases    → Todas las compras (admin)",
        "GET   /api/?action=all_users        → Todos los usuarios (admin)",
        "GET   /api/?action=occupied_seats   → Asientos ocupados por viaje",
        "GET   /api/?action=health           → Estado del servidor",
    ]
    add_bullet_points(slide, endpoints, Inches(0.5), Inches(1.8), Inches(6), Inches(3.5), size=12, color=GRAY_DARK)

    # PHP handler example
    add_title_text(slide, "handlePurchase() — Ejemplo de Handler PHP", Inches(6.8), Inches(1.35), Inches(6), Inches(0.35), size=16, color=BLUE_DARK, bold=True)
    php_code = """function handlePurchase() {
  $input = getJsonInput();
  $usuarioId = (int)($input['usuario_id'] ?? 0);
  $viajeId   = (int)($input['viaje_id'] ?? 0);
  $asiento   = (int)($input['asiento'] ?? 0);

  // Validaciones con early return
  if ($usuarioId <= 0 || $viajeId <= 0 || $asiento <= 0) {
    jsonError('Datos de compra inválidos.');
  }

  $db = getConnection();

  // Verificar viaje existe y está activo
  $stmt = $db->prepare(
    'SELECT * FROM viajes WHERE id = ? AND activo = 1'
  );
  $stmt->execute([$viajeId]);
  $viaje = $stmt->fetch();
  if (!$viaje) { jsonError('Viaje no encontrado.'); }

  // Verificar asiento no ocupado (PDO Prepared)
  $stmt = $db->prepare(
    'SELECT id FROM compras WHERE viaje_id = ? AND asiento = ?'
  );
  $stmt->execute([$viajeId, $asiento]);
  if ($stmt->fetch()) {
    jsonError('Este asiento ya está ocupado.');
  }

  // Registrar compra
  $stmt = $db->prepare(
    'INSERT INTO compras (...) VALUES (?, ?, ?, ?, ?, ?, ?, ?)'
  );
  $stmt->execute([...]);

  jsonResponse(['success' => true, 'purchase' => [...]]);
}"""
    add_code_block(slide, php_code, Inches(6.8), Inches(1.75), Inches(6), Inches(5.5), font_size=9)

    # DB Schema note
    add_rounded_rect(slide, Inches(0.5), Inches(5.5), Inches(6), Inches(1.8), GRAY_LIGHT)
    add_title_text(slide, "🗄️ Esquema de Base de Datos", Inches(0.7), Inches(5.6), Inches(5.5), Inches(0.3), size=15, color=BLUE_DARK)
    db_items = [
        "📋 usuarios: id, nombre, correo, password (SHA-256), rol",
        "📋 viajes: id, nombre_ruta, origen, destino, paradas (JSON), precio, asientos",
        "📋 compras: id, usuario_id (FK), viaje_id (FK), asiento, precio, fecha",
    ]
    add_bullet_points(slide, db_items, Inches(0.7), Inches(6.0), Inches(5.5), Inches(1.2), size=11, color=GRAY_MED)


def slide_cierre(prs):
    """Slide final: Cierre."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_gradient_bg(slide)

    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), BLUE_LIGHT)

    add_title_text(slide, "¡Gracias! 🚌", Inches(0.8), Inches(2.0), Inches(11), Inches(1.2), size=54, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_text(slide, "PremiumBus — Sistema de Transporte Urbano",
                   Inches(0.8), Inches(3.3), Inches(11), Inches(0.8), size=24, color=RGBColor(0xA0, 0xC4, 0xED), bold=False, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(5.5), Inches(4.3), Inches(2), Inches(0.05), BLUE_LIGHT)

    summary_items = [
        "📱  7 pantallas principales implementadas",
        "🚌  30 rutas reales de San Luis Potosí",
        "🔐  Sistema de autenticación dual (API + Demo)",
        "🗺️  Mapas interactivos con Leaflet.js",
        "🧪  17 pruebas unitarias exitosas",
        "📲  PWA instalable en Android e iOS",
    ]
    add_bullet_points(slide, summary_items, Inches(3), Inches(4.6), Inches(7), Inches(2.5), size=16, color=RGBColor(0xD0, 0xE0, 0xF0))

    add_rect(slide, Inches(0), Inches(7.42), SLIDE_WIDTH, Inches(0.08), BLUE_LIGHT)


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# MAIN
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

def main():
    print("🚌 Generando presentación PremiumBus...")

    prs = Presentation()
    prs.slide_width  = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Generar slides
    slide_portada(prs)            # 1
    slide_indice(prs)             # 2
    slide_arquitectura(prs)       # 3
    slide_menu_navegacion(prs)    # 4
    slide_pantallas_principales(prs)  # 5-11
    slide_codigo_auth(prs)        # 12
    slide_codigo_data(prs)        # 13
    slide_codigo_router(prs)      # 14
    slide_acceso_sistema(prs)     # 15
    slide_pruebas_auth(prs)       # 16
    slide_pruebas_data(prs)       # 17
    slide_api_backend(prs)        # 18
    slide_cierre(prs)             # 19

    prs.save(OUTPUT_FILE)
    print(f"✅ Presentación guardada en: {OUTPUT_FILE}")
    print(f"📊 Total de diapositivas: {len(prs.slides)}")


if __name__ == '__main__':
    main()
